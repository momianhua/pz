#!/usr/bin/env python3
"""Create isolated, valid Office fixtures from the lightweight contest samples."""

from __future__ import annotations

import argparse
import csv
import shutil
from pathlib import Path


def require_libraries():
    try:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit(
            "Missing test fixture dependency. Run: python -m pip install -r requirements-test.txt"
        ) from exc
    return Document, Workbook, Presentation, Inches, Pt


def build_report_docx(source: Path, target: Path, Document) -> None:
    text = source.read_text(encoding="utf-8-sig")
    document = Document()
    for index, raw in enumerate(text.splitlines()):
        line = raw.strip()
        if not line:
            continue
        if index == 0:
            document.add_heading(line, level=0)
        elif line == "执行摘要" or (line[:1].isdigit() and "." in line[:4]):
            document.add_heading(line, level=1 if line == "执行摘要" or line.count(".") == 1 else 2)
        else:
            document.add_paragraph(line)
    document.save(target)


def build_phone_docx(target: Path, Document) -> None:
    document = Document()
    document.add_heading("华为2025手机产品概览", level=0)
    document.add_paragraph("本文档介绍华为2025年主要手机产品线及其规格参数。")
    groups = {
        "Mate 系列": [
            ("Mate 70 Pro", "6.9英寸 OLED", "麒麟9020", 6499),
            ("Mate 70 Pro+", "6.9英寸 OLED", "麒麟9020", 8499),
            ("Mate 70 RS 非凡大师", "6.9英寸 OLED", "麒麟9020", 11999),
            ("Mate X6 折叠", "7.8英寸 OLED", "麒麟9020", 12999),
        ],
        "Pura 系列": [
            ("Pura 80", "6.7英寸 OLED", "麒麟9010", 4999),
            ("Pura 80 Pro", "6.8英寸 OLED", "麒麟9020", 6499),
            ("Pura 80 Ultra", "6.8英寸 OLED", "麒麟9020", 9999),
        ],
        "Nova 系列": [
            ("Nova 14", "6.7英寸 OLED", "麒麟8000", 2699),
            ("Nova 14 Pro", "6.78英寸 OLED", "麒麟9010", 3699),
            ("Nova 14 Ultra", "6.78英寸 OLED", "麒麟9010", 4699),
        ],
        "畅享 系列": [
            ("畅享 80", "6.67英寸 LCD", "天玑6020", 1299),
            ("畅享 80 Pro", "6.7英寸 OLED", "天玑7050", 1799),
        ],
    }
    for heading, rows in groups.items():
        document.add_heading(heading, level=1)
        table = document.add_table(rows=1, cols=4)
        for cell, value in zip(table.rows[0].cells, ["型号", "屏幕", "处理器", "起售价(元)"]):
            cell.text = value
        for row in rows:
            cells = table.add_row().cells
            for cell, value in zip(cells, row):
                cell.text = str(value)
    document.save(target)


def build_inventory_xlsx(source: Path, target: Path, Workbook) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "库存管理台账"
    with source.open("r", encoding="utf-8-sig", newline="") as stream:
        for row in csv.reader(stream, delimiter="\t"):
            sheet.append(row)
    workbook.save(target)


def build_video_pptx(target: Path, Presentation, Inches, Pt) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slides = [
        ("短视频平台差异化分析", ["测试用模拟数据集", "覆盖行业、用户、内容与推荐机制"]),
        ("核心结论", ["平台分化来自用户结构和推荐效率", "数据口径：2025 年模拟样本"]),
        ("行业概览", ["短视频用户渗透率 82%", "行业月活同比增长 9%"]),
        ("用户规模", ["平台 A 月活 7.2 亿", "平台 B 月活 4.1 亿"]),
        ("用户特征", ["平台 A 三线及以下用户占 58%", "平台 B 18—30 岁用户占 61%"]),
        ("内容生态", ["生活内容互动率 12.4%", "知识内容完播率 47%"]),
        ("推荐机制", ["兴趣召回贡献 64% 播放量", "冷启动平均耗时下降 18%"]),
    ]
    for title, bullets in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(24)
    presentation.save(target)


def build_delete_tree(target: Path) -> None:
    (target / "子目录").mkdir(parents=True, exist_ok=True)
    (target / "西安会议纪要.txt").write_text("应删除", encoding="utf-8")
    (target / "子目录" / "客户_西安_清单.csv").write_text("应删除", encoding="utf-8")
    (target / "北京会议纪要.txt").write_text("应保留", encoding="utf-8")
    (target / "子目录" / "西北区域.txt").write_text("应保留", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()
    source = args.source.resolve()
    output = args.output.resolve()
    if output.exists():
        raise SystemExit(f"Output already exists: {output}")
    shutil.copytree(source, output)

    Document, Workbook, Presentation, Inches, Pt = require_libraries()
    build_report_docx(
        source / "office_11" / "OpenClaw学术洞察报告.docx",
        output / "office_11" / "OpenClaw学术洞察报告.docx",
        Document,
    )
    build_phone_docx(output / "office_132" / "华为2025手机.docx", Document)
    build_inventory_xlsx(
        source / "office_18" / "generate_excel_.csv",
        output / "office_18" / "generate_excel_1.xlsx",
        Workbook,
    )
    build_video_pptx(output / "office_022" / "短视频平台差异化分析报告.pptx", Presentation, Inches, Pt)
    build_delete_tree(output / "office_103" / "待清理目录")
    print(output)


if __name__ == "__main__":
    main()
